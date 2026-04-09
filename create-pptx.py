#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Colour palette (Dynatrace-inspired) ──
DT_PURPLE = RGBColor(0x61, 0x2C, 0xA6)
DT_DARK = RGBColor(0x1A, 0x1A, 0x2E)
DT_GREEN = RGBColor(0x14, 0xBE, 0xA0)
DT_BLUE = RGBColor(0x00, 0x84, 0xBF)
DT_ORANGE = RGBColor(0xF5, 0x89, 0x28)
DT_RED = RGBColor(0xE8, 0x45, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xE0, 0xE0, 0xE0)
MID_GREY = RGBColor(0xA0, 0xA0, 0xA0)
DARK_TEXT = RGBColor(0x2A, 0x2A, 0x3C)

def add_bg(slide, colour):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colour

def add_shape(slide, left, top, width, height, fill_colour, border_colour=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_colour
    if border_colour:
        shape.line.color.rgb = border_colour
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, colour=WHITE, align=PP_ALIGN.LEFT, font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = colour
    p.font.name = font_name
    p.alignment = align
    return tf

def add_bullet_list(tf, items, font_size=13, colour=WHITE, bold_first=False, indent=0, font_name='Segoe UI'):
    for i, item in enumerate(items):
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = colour
        p.font.name = font_name
        p.space_before = Pt(4)
        p.level = indent
        if bold_first and i == 0:
            p.font.bold = True

# ═════════════════════════════════════════════════════════════
# SLIDE 1 — BUSINESS OBSERVABILITY DEMONSTRATOR: POWER & WINS
# ═════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide1, DT_DARK)

# ── Title bar ──
add_shape(slide1, Inches(0), Inches(0), Inches(13.333), Inches(1.15), DT_PURPLE)
add_text_box(slide1, Inches(0.6), Inches(0.15), Inches(10), Inches(0.55),
             'Business Observability Demonstrator', font_size=28, bold=True, colour=WHITE)
add_text_box(slide1, Inches(0.6), Inches(0.65), Inches(10), Inches(0.4),
             'Wins, Partners & Opportunities — Q2 2026', font_size=15, colour=LIGHT_GREY)

# ── Dynatrace logo text (top-right) ──
add_text_box(slide1, Inches(11.0), Inches(0.25), Inches(2), Inches(0.5),
             'dynatrace', font_size=22, bold=True, colour=DT_GREEN, align=PP_ALIGN.RIGHT)

# ── LEFT COLUMN: Wins This Quarter ──
card1 = add_shape(slide1, Inches(0.5), Inches(1.45), Inches(4.0), Inches(2.6), RGBColor(0x25, 0x25, 0x40), DT_GREEN)
tf1 = add_text_box(slide1, Inches(0.8), Inches(1.55), Inches(3.5), Inches(0.4),
                   '🏆  WINS THIS QUARTER', font_size=16, bold=True, colour=DT_GREEN)
add_bullet_list(tf1, [
    'Neso — Partner: Wipro',
    'Met Police',
    'Parkland Health — Partner: Kyndryl',
], font_size=14, colour=WHITE)
# Add emphasis line
tf1b = add_text_box(slide1, Inches(0.8), Inches(3.35), Inches(3.5), Inches(0.5),
                    'Deals rescued by tying IT back to business outcomes', font_size=11, colour=MID_GREY)

# ── MIDDLE COLUMN: Partners Deployed ──
card2 = add_shape(slide1, Inches(4.8), Inches(1.45), Inches(4.0), Inches(2.6), RGBColor(0x25, 0x25, 0x40), DT_BLUE)
tf2 = add_text_box(slide1, Inches(5.1), Inches(1.55), Inches(3.5), Inches(0.4),
                   '🤝  PARTNERS DEPLOYED', font_size=16, bold=True, colour=DT_BLUE)
add_bullet_list(tf2, [
    'Accenture',
    'TCS',
    'Eviden',
    'Alanata',
    'Home',
    'TietoEvry',
], font_size=14, colour=WHITE)

# ── RIGHT COLUMN: Key Account Spotlight ──
card3 = add_shape(slide1, Inches(9.1), Inches(1.45), Inches(3.8), Inches(2.6), RGBColor(0x25, 0x25, 0x40), DT_ORANGE)
tf3 = add_text_box(slide1, Inches(9.4), Inches(1.55), Inches(3.3), Inches(0.4),
                   '🔑  KEY ACCOUNT: MUFG', font_size=16, bold=True, colour=DT_ORANGE)
add_bullet_list(tf3, [
    'Mitsubishi UFJ Financial Group',
    'Uses Murex MX.3 for trading, risk & post-trade ops',
    'EMEA securities + Krungsri subsidiary',
    'BizObs ties trading platform health → revenue impact',
], font_size=12, colour=WHITE)

# ── BOTTOM ROW: Why Business Observability Matters ──
card4 = add_shape(slide1, Inches(0.5), Inches(4.35), Inches(12.4), Inches(2.8), RGBColor(0x20, 0x20, 0x38), DT_PURPLE)

tf4 = add_text_box(slide1, Inches(0.9), Inches(4.45), Inches(5.5), Inches(0.4),
                   '📊  WHY BUSINESS OBSERVABILITY MATTERS', font_size=16, bold=True, colour=DT_GREEN)

# Left sub-column
tf4a = add_text_box(slide1, Inches(0.9), Inches(5.0), Inches(5.5), Inches(2.0),
                    '', font_size=13, colour=WHITE)
tf4a.paragraphs[0].text = 'The Problem'
tf4a.paragraphs[0].font.bold = True
tf4a.paragraphs[0].font.size = Pt(14)
tf4a.paragraphs[0].font.color.rgb = DT_RED
add_bullet_list(tf4a, [
    'Deals at risk when IT metrics can\'t tie back to business outcomes',
    'Executives don\'t care about p99 latency — they care about revenue',
    'Competitors (Leonardo, etc.) reframe the narrative around business value',
    'Without BizObs, we\'re selling infrastructure monitoring — not transformation',
], font_size=12, colour=LIGHT_GREY)

# Right sub-column
tf4b = add_text_box(slide1, Inches(7.0), Inches(5.0), Inches(5.5), Inches(2.0),
                    '', font_size=13, colour=WHITE)
tf4b.paragraphs[0].text = 'The Demonstrator Advantage'
tf4b.paragraphs[0].font.bold = True
tf4b.paragraphs[0].font.size = Pt(14)
tf4b.paragraphs[0].font.color.rgb = DT_GREEN
add_bullet_list(tf4b, [
    'Not taking systems out — moving the narrative to business outcomes',
    'Live, tailored demo in 30 min: their journey, their KPIs, their revenue',
    '110+ industry templates across 55+ verticals',
    'AI-powered: describe any journey in plain language → full simulation',
    'Every failure shows business impact: "340 journeys blocked, $127K at risk"',
], font_size=12, colour=LIGHT_GREY)


# ═════════════════════════════════════════════════════════════
# SLIDE 2 — THE PLAYBOOK: HOW TO WIN WITH BIZOBS
# ═════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide2, DT_DARK)

# ── Title bar ──
add_shape(slide2, Inches(0), Inches(0), Inches(13.333), Inches(1.15), DT_PURPLE)
add_text_box(slide2, Inches(0.6), Inches(0.15), Inches(10), Inches(0.55),
             'The Playbook: Deploy → Demo → Win', font_size=28, bold=True, colour=WHITE)
add_text_box(slide2, Inches(0.6), Inches(0.65), Inches(10), Inches(0.4),
             'Simple deployment, powerful story, proven results', font_size=15, colour=LIGHT_GREY)
add_text_box(slide2, Inches(11.0), Inches(0.25), Inches(2), Inches(0.5),
             'dynatrace', font_size=22, bold=True, colour=DT_GREEN, align=PP_ALIGN.RIGHT)

# ── LEFT: The Story Structure ──
card5 = add_shape(slide2, Inches(0.5), Inches(1.45), Inches(6.2), Inches(3.1), RGBColor(0x25, 0x25, 0x40), DT_GREEN)

tf5 = add_text_box(slide2, Inches(0.8), Inches(1.55), Inches(5.8), Inches(0.4),
                   '🎯  THE STORY (Tell → Show → Close)', font_size=16, bold=True, colour=DT_GREEN)

# Step 1
tf5a = add_text_box(slide2, Inches(0.8), Inches(2.05), Inches(5.6), Inches(0.35),
                    '1.  Tell \'em what you\'re gonna tell them', font_size=14, bold=True, colour=DT_ORANGE)
tf5a2 = add_text_box(slide2, Inches(1.2), Inches(2.35), Inches(5.2), Inches(0.35),
                     '"Your deals are at risk. IT metrics alone don\'t close executives."', font_size=12, colour=LIGHT_GREY)

# Step 2
tf5b = add_text_box(slide2, Inches(0.8), Inches(2.75), Inches(5.6), Inches(0.35),
                    '2.  Tell them — Live Demo', font_size=14, bold=True, colour=DT_ORANGE)
tf5b2 = add_text_box(slide2, Inches(1.2), Inches(3.05), Inches(5.2), Inches(0.6),
                     '', font_size=12, colour=LIGHT_GREY)
tf5b2.paragraphs[0].text = 'Show their actual journey running live in Dynatrace'
tf5b2.paragraphs[0].font.size = Pt(12)
tf5b2.paragraphs[0].font.color.rgb = LIGHT_GREY
add_bullet_list(tf5b2, [
    'Break something → show revenue impact, not just errors',
    'AI resolves it → show autonomous operations closing the loop',
], font_size=11, colour=MID_GREY)

# Step 3
tf5c = add_text_box(slide2, Inches(0.8), Inches(3.75), Inches(5.6), Inches(0.35),
                    '3.  Tell \'em what you told them', font_size=14, bold=True, colour=DT_ORANGE)
tf5c2 = add_text_box(slide2, Inches(1.2), Inches(4.05), Inches(5.2), Inches(0.35),
                     '"Get it deployed → I\'ll help with enablement post-deployment."', font_size=12, colour=LIGHT_GREY)

# ── RIGHT: Simple Deployment Steps ──
card6 = add_shape(slide2, Inches(7.0), Inches(1.45), Inches(5.9), Inches(3.1), RGBColor(0x25, 0x25, 0x40), DT_BLUE)

tf6 = add_text_box(slide2, Inches(7.3), Inches(1.55), Inches(5.3), Inches(0.4),
                   '🚀  DEPLOYMENT (< 30 minutes)', font_size=16, bold=True, colour=DT_BLUE)

steps = [
    ('①', 'Run setup script', 'Installs Node.js, Ollama AI, OneAgent, OpenTelemetry'),
    ('②', 'Deploy to Dynatrace', 'One command: npx dt-app deploy → live in your tenant'),
    ('③', 'Configure EdgeConnect', 'Secure tunnel from Dynatrace to your server'),
    ('④', 'Open the app → Get Started', '5-tab wizard walks through everything'),
    ('⑤', 'Choose a journey → Go', 'Pick from 110+ templates or describe your own'),
]

y_pos = 2.05
for num, title, desc in steps:
    tf_step = add_text_box(slide2, Inches(7.3), Inches(y_pos), Inches(5.3), Inches(0.22),
                           f'{num}  {title}', font_size=13, bold=True, colour=WHITE)
    tf_desc = add_text_box(slide2, Inches(7.8), Inches(y_pos + 0.22), Inches(4.8), Inches(0.22),
                           desc, font_size=11, colour=MID_GREY)
    y_pos += 0.52

# ── BOTTOM: Results & CTA ──
card7 = add_shape(slide2, Inches(0.5), Inches(4.85), Inches(12.4), Inches(2.3), RGBColor(0x20, 0x20, 0x38), DT_PURPLE)

# Results header
tf7 = add_text_box(slide2, Inches(0.9), Inches(4.95), Inches(5), Inches(0.4),
                   '📈  PROVEN RESULTS', font_size=16, bold=True, colour=DT_GREEN)

# Results grid
results = [
    ('Neso + Wipro', 'BizObs demo shifted the conversation\nfrom infra monitoring to business outcomes'),
    ('Parkland Health + Kyndryl', 'Patient journey simulation showed\nexecutive team real-time revenue at risk'),
    ('Met Police', 'Connected operational systems\nto citizen-facing service outcomes'),
]

x_pos = 0.9
for title, desc in results:
    tf_r_title = add_text_box(slide2, Inches(x_pos), Inches(5.4), Inches(3.7), Inches(0.3),
                              title, font_size=13, bold=True, colour=DT_ORANGE)
    tf_r_desc = add_text_box(slide2, Inches(x_pos), Inches(5.7), Inches(3.7), Inches(0.6),
                             desc, font_size=11, colour=LIGHT_GREY)
    x_pos += 4.1

# CTA
tf_cta = add_text_box(slide2, Inches(0.9), Inches(6.55), Inches(11.5), Inches(0.4),
                      '💬  Get it deployed → I can help with enablement and partner onboarding.   |   Not taking systems out — moving the narrative to Business Observability.',
                      font_size=13, bold=True, colour=DT_GREEN)

# ── Save ──
output_path = '/home/ec2-user/Dynatrace-Business-Observability-Forge/Business-Observability-Demonstrator-Wins.pptx'
prs.save(output_path)
print(f'✅ PowerPoint saved to: {output_path}')
